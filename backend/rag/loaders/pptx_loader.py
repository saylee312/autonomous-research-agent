from pptx import Presentation


def load_pptx(file_path):

    prs = Presentation(file_path)

    result = {

        "text": [],

        "tables": [],

        "images": []
    }

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(
                shape,
                "text"
            ):

                if shape.text:

                    result["text"].append(
                        shape.text
                    )

    return result